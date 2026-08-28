#!/usr/bin/env python3
"""
build_showcase_pptx.py — generate docs/Tech_Drawdown_Executive_Showcase.pptx

The executive deck is now GENERATED, like the BRD / Architecture / Playbook docs,
so it stays in sync with the shipped feature set instead of drifting as a
hand-maintained artifact.

Design system (extracted from the original hand-built deck):
  · 16:9, 13.33in x 7.5in
  · Display font Cambria, body font Calibri
  · Light slides alternate FFFFFF / F4F6F9; hero + method + value + closing are navy 0F2440
  · Eyebrow (small caps, colored) -> Title (Cambria bold) -> content -> footer + page no.
  · Cards: white fill, E5E9F0 hairline. Dark cards: 14304F fill, 1E3E5F hairline.

Icons live in docs/deck_assets/ (extracted from the original deck).
Diagrams are pulled straight from docs/diagrams/.

Usage:  python scripts/build_showcase_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
DOCS = os.path.join(BASE, "docs")
ASSETS = os.path.join(DOCS, "deck_assets")
DIAGRAMS = os.path.join(DOCS, "diagrams")
OUT = os.path.join(DOCS, "Tech_Drawdown_Executive_Showcase.pptx")

VERSION = "1.0"
DATE_LABEL = "AUGUST 2026"

# ---------------------------------------------------------------- design tokens
NAVY   = "0F2440"
DARKC  = "14304F"
DARKL  = "1E3E5F"
GOLD   = "E0A106"
GREEN  = "1D9E75"
AMBER  = "E0A93B"
RED    = "E24B4A"
BLUE   = "2E75B6"
PURPLE = "7C3AED"
TEAL   = "0E7C86"
INK    = "1F2937"
MUTED  = "64748B"
LINE   = "E5E9F0"
CARD   = "FFFFFF"
SOFT   = "F4F6F9"
LIGHT  = "CBD5E1"
SLATE  = "94A3B8"
STEEL  = "7C93AB"
LILAC  = "F3EEFC"

DISPLAY, BODY = "Cambria", "Calibri"

IN = 914400
M = 640080                      # left margin
CONTENT_W = 10908792            # content width
FOOTER_Y = 6492240
SLIDE_W, SLIDE_H = 12188952, 6858000


def C(h):
    return RGBColor.from_string(h)


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W)
        self.prs.slide_height = Emu(SLIDE_H)
        self.n = 0

    # -- primitives ------------------------------------------------------------
    def slide(self, bg):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.n += 1
        # background
        from lxml import etree
        ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        bge = etree.SubElement(s._element, f"{{{ns}}}bg")
        pr = etree.SubElement(bge, f"{{{ns}}}bgPr")
        fill = etree.SubElement(pr, f"{{{a}}}solidFill")
        clr = etree.SubElement(fill, f"{{{a}}}srgbClr")
        clr.set("val", bg)
        etree.SubElement(pr, f"{{{a}}}effectLst")
        s._element.remove(bge)
        s._element.insert(0, bge)
        return s

    def text(self, s, x, y, w, h, runs, size=18, bold=False, color=INK,
             font=BODY, align=PP_ALIGN.LEFT, space=0):
        """runs: str, or list of (text, {opts}) tuples, or list of lines (list of runs)."""
        tb = s.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        if isinstance(runs, str):
            lines = [[(runs, {})]]
        elif runs and isinstance(runs[0], tuple):
            lines = [runs]
        else:
            lines = runs
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if space:
                p.space_after = Pt(space)
            for txt, o in ln:
                r = p.add_run()
                r.text = txt
                f = r.font
                f.size = Pt(o.get("size", size))
                f.bold = o.get("bold", bold)
                f.name = o.get("font", font)
                f.color.rgb = C(o.get("color", color))
        return tb

    def rect(self, s, x, y, w, h, fill=None, line=None, lw=1, radius=None):
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
            Emu(x), Emu(y), Emu(w), Emu(h))
        if radius:
            try:
                shp.adjustments[0] = radius
            except Exception:
                pass
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = C(fill)
        else:
            shp.fill.background()
        if line:
            shp.line.color.rgb = C(line)
            shp.line.width = Pt(lw)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def icon(self, s, x, y, box, color, img):
        """Colored rounded tile with a centered white icon."""
        self.rect(s, x, y, box, box, fill=color, radius=0.18)
        path = os.path.join(ASSETS, img)
        if os.path.exists(path):
            sz = int(box * 0.42)
            s.shapes.add_picture(path, Emu(x + (box - sz) // 2), Emu(y + (box - sz) // 2),
                                 Emu(sz), Emu(sz))

    def chrome(self, s, eyebrow, title, eyecolor=GOLD, dark=False, sub=None):
        tcol = CARD if dark else INK
        self.text(s, M, 457200, 8229600, 274320, eyebrow, size=12, bold=True,
                  color=eyecolor, font=BODY)
        self.text(s, M, 566928, CONTENT_W, 731520, title, size=32, bold=True,
                  color=tcol, font=DISPLAY)
        if sub:
            self.text(s, M, 1371600, CONTENT_W, 548640, sub, size=14.5,
                      color=LIGHT if dark else MUTED)
        self.footer(s, dark)

    def footer(self, s, dark=False):
        col = STEEL if dark else MUTED
        self.text(s, M, FOOTER_Y, 5486400, 274320,
                  "Tech Drawdown · Executive Showcase", size=9, color=col)
        self.text(s, 11091672, FOOTER_Y, 457200, 274320, str(self.n),
                  size=9, color=col, align=PP_ALIGN.RIGHT)

    # -- composite layouts -----------------------------------------------------
    def kpis(self, s, items, top=1965960):
        """items: [(big, bigcolor, caption)] — 3 or 4 up."""
        n = len(items)
        gap = 274320
        w = (CONTENT_W - gap * (n - 1)) // n
        for i, (big, col, cap) in enumerate(items):
            x = M + i * (w + gap)
            self.rect(s, x, top, w, 2286000, fill=CARD, line=LINE, radius=0.06)
            self.text(s, x + 274320, top + 320040, w - 548640, 914400, big,
                      size=40, bold=True, color=col, font=DISPLAY)
            self.text(s, x + 274320, top + 1280160, w - 548640, 822960, cap,
                      size=13.5, color=MUTED)

    def tiles(self, s, items, top=1783080, h=3291840, dark=False):
        """items: [(icon_png, icon_color, heading, body)] — 2, 3 or 4 up."""
        n = len(items)
        gap = 274320
        w = (CONTENT_W - gap * (n - 1)) // n
        for i, (img, col, head, body) in enumerate(items):
            x = M + i * (w + gap)
            self.rect(s, x, top, w, h,
                      fill=DARKC if dark else CARD,
                      line=DARKL if dark else LINE, radius=0.05)
            px, py = x + 256032, top + 256032
            self.icon(s, px, py, 603504, col, img)
            self.text(s, px, top + 960120, w - 512064, 365760, head,
                      size=15, bold=True, color=CARD if dark else INK)
            self.text(s, px, top + 1325880, w - 512064, h - 1600200, body,
                      size=12, color=LIGHT if dark else MUTED)

    def specrows(self, s, rows, top=1783080):
        """rows: [(label, value)] — banded two-column grid."""
        rh = 658368
        for i, (lab, val) in enumerate(rows):
            y = top + i * rh
            if i % 2 == 0:
                self.rect(s, M, y, CONTENT_W, rh * 2 if i + 1 < len(rows) else rh,
                          fill=SOFT, line=SOFT)
        for i, (lab, val) in enumerate(rows):
            y = top + i * rh
            self.text(s, M + 182880, y + 137160, 2468880, 400000, lab,
                      size=14, bold=True, color=BLUE)
            self.text(s, 3383280, y + 137160, 7982712, 400000, val,
                      size=13.5, color=INK)

    def bullets(self, s, items, x=M, y=1828800, w=6126480, h=2926080, size=15, dark=False):
        lines = [[("•  " + t, {})] for t in items]
        self.text(s, x, y, w, h, lines, size=size,
                  color=LIGHT if dark else INK, space=10)

    def image(self, s, name, x, y, w, h):
        path = name if os.path.isabs(name) else os.path.join(DIAGRAMS, name)
        if not os.path.exists(path):
            return
        from PIL import Image
        iw, ih = Image.open(path).size
        sc = min(w / iw, h / ih)
        dw, dh = int(iw * sc), int(ih * sc)
        s.shapes.add_picture(path, Emu(x + (w - dw) // 2), Emu(y + (h - dh) // 2),
                             Emu(dw), Emu(dh))

    def save(self):
        self.prs.save(OUT)
        print("wrote", OUT, os.path.getsize(OUT), "bytes,", self.n, "slides")


d = Deck()

# ── 1 · HERO ────────────────────────────────────────────────────────────────────
s = d.slide(NAVY)
for i, col in enumerate((GREEN, AMBER, RED)):
    d.rect(s, M + i * 384048, 2148840, 310896, 310896, fill=col, radius=0.2)
d.text(s, 658368, 1554480, 8229600, 274320, f"EXECUTIVE SHOWCASE · {DATE_LABEL}",
       size=12, bold=True, color=GOLD)
d.text(s, 603504, 2514600, 10515600, 1097280, "Tech Drawdown",
       size=54, bold=True, color=CARD, font=DISPLAY)
d.text(s, 658368, 3703320, 9692640, 1005840,
       "A live market-drawdown dashboard, private-portfolio tracker & retirement planner — "
       "designed, built, hardened, and documented through conversation with an AI agent.",
       size=17, color=LIGHT)
d.text(s, 658368, 5029200, 9144000, 365760, [
    (f"Version {VERSION}", {"bold": True, "color": GOLD}),
    ("   ·   Rules-based analytics, not financial advice", {"color": SLATE}),
], size=13)

# ── 2 · WHY THIS MATTERS ────────────────────────────────────────────────────────
s = d.slide(SOFT)
d.text(s, M, 457200, 8229600, 274320, "WHY THIS MATTERS FOR AI ADOPTION",
       size=12, bold=True, color=GOLD)
d.text(s, M, 777240, CONTENT_W, 731520,
       "One real, production-grade application — built end to end in conversation.",
       size=26, bold=True, color=INK, font=DISPLAY)
d.kpis(s, [
    ("~1,100", BLUE,  "securities monitored across six market universes"),
    ("Zero-touch", GREEN, "daily refresh — fully unattended automation"),
    ("1 file", GOLD,  "self-contained web page; opens anywhere, nothing to install"),
])
d.text(s, M, 4709160, CONTENT_W, 822960,
       "No dev team spun up. No infrastructure stood up. A spreadsheet-and-manual-lookup "
       "problem became a self-updating system — the same pattern you can bring to client workflows.",
       size=15, color=INK)
d.footer(s)

# ── 3 · THE PROBLEM ─────────────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "THE PROBLEM IT REPLACES", "Spreadsheets couldn't keep up", eyecolor=RED)
d.tiles(s, [
    ("image-3-1.png", RED, "Stale, error-prone data",
     "Values and drawdowns were only as fresh as the last manual refresh; off-index holdings "
     "were mis-valued from outdated figures."),
    ("image-3-2.png", RED, "No single view",
     "Market context, personal holdings, and retirement projections lived in separate files — "
     "hard to see the whole picture or act quickly."),
    ("image-3-3.png", RED, "High manual effort",
     "Keeping everything current consumed hours every week and was easy to forget or get wrong."),
], top=1783080, h=2926080)
d.text(s, M, 5029200, 5486400, 640080,
       [[("Hours of upkeep,", {})], [("every week.", {})]],
       size=20, bold=True, color=INK, font=DISPLAY)
d.text(s, M, 5760720, 6858000, 365760,
       "Manual, fragmented, and always a little out of date.", size=13, color=MUTED)

# ── 4 · THE SOLUTION ────────────────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "THE SOLUTION", "One always-current view", eyecolor=GREEN)
d.text(s, M, 1828800, 5852160, 2926080, [
    [("A single, self-updating web page that unifies three things at a glance:", {"bold": True})],
    [("•  Which stocks and funds are in a drawdown — and how deep", {})],
    [("•  What every holding is worth at today's live prices", {})],
    [("•  Whether the household is on track for retirement income", {})],
    [("…and now: how it performs against the market, what changed overnight, "
      "and what it will pay in dividends.", {"color": MUTED})],
], size=16, color=INK, space=10)
quad = [("image-4-1.png", BLUE,  "Market drawdown & trend"),
        ("image-4-2.png", GREEN, "Live private portfolios"),
        ("image-4-3.png", GOLD,  "Retirement & income"),
        ("image-4-4.png", PURPLE, "Fundamentals + factsheet")]
for i, (img, col, lab) in enumerate(quad):
    cx = 6949440 + (i % 2) * 2468880
    cy = 1828800 + (i // 2) * 1417320
    d.rect(s, cx, cy, 2286000, 1234440, fill=CARD, line=LINE, radius=0.06)
    d.icon(s, cx + 201168, cy + 219456, 502920, col, img)
    d.text(s, cx + 182880, cy + 822960, 1965960, 365760, lab, size=12.5, bold=True, color=INK)

# ── 5 · AT A GLANCE ─────────────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "AT A GLANCE", "What it delivers")
d.specrows(s, [
    ("Coverage", "~1,100 securities — S&P 500, Nasdaq-100, Dow, Top-100 ETFs, Thematic ETFs, Mutual Funds"),
    ("Private accounts", "IRA (income focus) + Brokerage (growth focus) + a household retirement planner"),
    ("Refresh", "Automatic daily at 9am, plus an on-demand 'Load live prices' button"),
    ("History & income", "Rolling daily value history → benchmark vs Dow/Nasdaq/S&P, drawdown timeline, "
                         "change alerts, and a forward 12-month dividend calendar"),
    ("Delivery", "Single self-contained web page — opens in any browser, works offline"),
    ("Data privacy", "Holdings stay local; only ticker symbols are sent, only for public market data"),
    ("Reliability", "Integrity-gated automation — always falls back to the last good state, never blanks a holding"),
], top=1600200)

# ── 6 · OVERVIEW SPECTRUM ───────────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "FEATURE · MARKET-CONTEXT OVERVIEW", "A health spectrum, at a glance",
         eyecolor=GREEN,
         sub="The whole universe, de-duplicated and sorted into three color-coded columns — "
             "every ticker tagged with the index it belongs to.")
cols = [
    ("Best performing", GREEN, [("MRNA", "S&P", 1.0), ("WCBR", "THM", 0.72), ("IVE", "ETF", 0.44)]),
    ("In correction", AMBER,  [("MU", "S&P·NDX", 1.0), ("MCD", "S&P·DOW", 0.72), ("PBW", "THM", 0.44)]),
    ("Worst drawdown", RED,   [("TTD", "S&P·NDX", 1.0), ("FISV", "S&P", 0.72), ("MSTR", "NDX", 0.44)]),
]
for ci, (head, col, rows) in enumerate(cols):
    x = M + ci * 3703320
    d.rect(s, x, 2103120, 3429000, 3566160, fill=CARD, line=LINE, radius=0.05)
    d.rect(s, x, 2103120, 3429000, 548640, fill=col)
    d.text(s, x + 182880, 2194560, 3063240, 365760, head, size=15, bold=True, color=CARD)
    for ri, (tk, tags, frac) in enumerate(rows):
        y = 2880360 + ri * 658368
        d.text(s, x + 274320, y, 1280160, 365760, tk, size=14, bold=True, color=INK)
        d.text(s, x + 1371600, y + 27432, 1828800, 320040, tags, size=11, color=MUTED)
        d.rect(s, x + 274320, y + 384048, int(2651760 * frac), 82296, fill=col)
    d.text(s, x + 274320, 5257800, 2926080, 274320, "+ top-10 · click any ticker",
           size=10.5, color=MUTED)

# ── 7 · 1-YEAR CHARTS ───────────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "FEATURE · MARKET CONTEXT", "One-year charts — indices & macro")
for i, (img, col, head, body) in enumerate([
    ("image-7-1.png", BLUE, "US indices · 1-year",
     "Own-scale sparklines for the Dow, Nasdaq, S&P 500 and Russell 2000 — read each index "
     "directly, not rescaled against the others. Built once, embedded, no runtime calls."),
    ("image-7-2.png", GOLD, "Macro & commodities · 1-year",
     "Oil, gold, silver, copper, rare earths (ETF proxy), the US dollar vs four major currencies, "
     "and consumer sentiment — the wider backdrop behind the portfolio, in one card. Any series "
     "the data plan gates is simply omitted; the card shows the rest."),
]):
    x = M + i * 5669280
    d.rect(s, x, 1737360, 5394960, 3566160, fill=CARD, line=LINE, radius=0.05)
    d.icon(s, x + 256032, 1993392, 603504, col, img)
    d.text(s, x + 256032, 2697480, 4882896, 365760, head, size=15, bold=True, color=INK)
    d.text(s, x + 256032, 3063240, 4882896, 2103120, body, size=12, color=MUTED)

# ── 8 · PRIVATE PORTFOLIOS ──────────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "FEATURE · PRIVATE PORTFOLIOS", "Your holdings, valued live", eyecolor=GREEN)
d.bullets(s, [
    "Market value = quantity × today's live price",
    "IRA (income) and Brokerage (growth) tabs, each with allocation and gain/loss",
    "A guarded pricing rule rejects a corrupt quote and falls back safely — one bad "
    "number can never mis-state the portfolio",
    "Off-index holdings priced individually, so nothing is valued from a stale figure",
], y=1828800, w=6126480, h=2926080, size=15)
d.rect(s, 7406640, 1828800, 4114800, 2377440, fill=CARD, line=LINE, radius=0.05)
d.icon(s, 7680960, 2057400, 603504, GREEN, "image-8-1.png")
d.text(s, 7680960, 2761488, 3566160, 1188720,
       "Edit a flat export — the next run re-imports it and prices new names automatically.",
       size=13, color=MUTED)
d.text(s, M, 5029200, 6126480, 365760,
       "Holdings never leave the machine. Only ticker symbols go out.",
       size=13, bold=True, color=GREEN)

# ── 9 · RETIREMENT & INCOME ─────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "FEATURE · RETIREMENT & INCOME", "Planned to a retirement date")
d.tiles(s, [
    ("image-9-1.png", GOLD, "Two-account model",
     "IRA modeled for income, Brokerage for growth — roles baked into the projection."),
    ("image-9-2.png", GOLD, "Projections",
     "Household balances and sustainable income projected to your target date."),
    ("image-9-3.png", GOLD, "Rebalance targets",
     "See current plan vs. a rebalance plan, with the difference called out."),
])

# ── 10 · DIVIDEND INCOME CALENDAR ──────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "FEATURE · INCOME", "What will it pay me, and when?", eyecolor=TEAL,
         sub="A forward 12-month dividend calendar across every holding — the question an "
             "income portfolio exists to answer.")
d.bullets(s, [
    "Projected annual income, forward yield on market value, and income due in the next 30 days",
    "A month-by-month chart split by account — IRA (income engine) vs Brokerage",
    "An upcoming-payments list with exact ex-dividend, record and pay dates",
], y=2103120, w=6126480, h=1828800, size=15)
d.rect(s, M, 4114800, 6126480, 1097280, fill=CARD, line=LINE, radius=0.05)
d.text(s, M + 228600, 4251960, 5669280, 822960, [
    [("Honest by design.  ", {"bold": True, "color": TEAL}),
     ("Declared payments show an exact date and render green. Where the data plan gates "
      "per-symbol dates (ETFs and funds), income is estimated from yield × value, placed at "
      "month level, and labelled ", {}),
     ("estimated", {"bold": True}), (" — never silently faked.", {})],
], size=12, color=INK)
# mock calendar card
d.rect(s, 7406640, 2103120, 4114800, 3108960, fill=CARD, line=LINE, radius=0.05)
d.text(s, 7680960, 2331720, 3566160, 365760, "Dividend income", size=14, bold=True, color=INK)
d.text(s, 7680960, 2697480, 3566160, 320040, "Projected next 12 months", size=10.5, color=MUTED)
bars = [0.35, 0.42, 0.95, 0.38, 0.40, 0.92, 0.36, 0.44, 0.90, 0.34, 0.41, 0.88]
bw, bx0, by0, bmax = 274320, 7680960, 4297680, 1005840
for i, f in enumerate(bars):
    bh = int(bmax * f)
    d.rect(s, bx0 + i * (bw + 27432), by0 - bh, bw, bh,
           fill=TEAL if i % 3 != 2 else BLUE)
d.text(s, 7680960, 4389120, 3566160, 274320,
       "quarter-end months stack on the monthly ETFs", size=9.5, color=MUTED)
d.rect(s, 7680960, 4754880, 3566160, 274320, fill=LILAC)
d.text(s, 7818120, 4772000, 3383280, 240000, [
    [("declared", {"bold": True, "color": GREEN}), ("  ·  exact date    ", {"color": MUTED}),
     ("estimated", {"bold": True, "color": AMBER}), ("  ·  month level", {"color": MUTED})],
], size=9.5)

# ── 11 · BENCHMARK + TIMELINE ───────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "FEATURE · PERFORMANCE", "Am I beating the market?", eyecolor=BLUE,
         sub="A rolling daily history turned a point-in-time snapshot into something that can "
             "answer 'versus what?' and 'since when?'.")
d.tiles(s, [
    ("image-14-1.png", BLUE, "Portfolio vs the market",
     "Household value indexed to 100 against the Dow, Nasdaq and S&P 500, with a live alpha "
     "readout. No fabricated back-cast — the comparison accrues honestly from the day it was "
     "switched on."),
    ("image-14-2.png", RED, "Drawdown timeline",
     "An underwater chart: how far the portfolio and each index sit below their own running "
     "peak over the past year, with the Correction (−10%) and Bear (−20%) bands shaded."),
], top=2103120, h=2926080)

# ── 12 · CHANGE ALERTS ──────────────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "FEATURE · SIGNAL", "What changed since yesterday?", eyecolor=AMBER,
         sub="Stop re-scanning the whole board every morning to find the two things that moved.")
d.bullets(s, [
    "A daily diff of held names: entered Bear or Correction, Golden or Death Cross, "
    "recovered, or hit a new 52-week low",
    "Shown as a colour-coded icon beside the ticker — click it for the from→to detail and history",
    "Listed in the automated run report, so the change is in the record too",
], y=2103120, w=6126480, h=1828800, size=15)
d.rect(s, M, 4206240, 6126480, 1005840, fill=CARD, line=LINE, radius=0.05)
d.text(s, M + 228600, 4343400, 5669280, 731520, [
    [("Zero new data calls.  ", {"bold": True, "color": AMBER}),
     ("Every field it needs is already on hand. And it is guarded by the integrity check — a "
      "corrupt price can never fabricate a Death Cross. An alert system that cries wolf once "
      "is worse than none at all.", {})],
], size=12, color=INK)
# mock alert rows
d.rect(s, 7406640, 2103120, 4114800, 3108960, fill=CARD, line=LINE, radius=0.05)
d.text(s, 7680960, 2331720, 3566160, 320040, "Changes since yesterday", size=13, bold=True, color=INK)
alerts = [("SNOW", "Golden Cross", GREEN), ("NVO", "Death Cross", RED),
          ("ACHR", "Death Cross", RED), ("ASML", "entered Correction", AMBER),
          ("META", "recovered to Correction", GREEN)]
for i, (tk, lab, col) in enumerate(alerts):
    y = 2789136 + i * 466344
    d.rect(s, 7680960, y + 45720, 137160, 137160, fill=col, radius=0.5)
    d.text(s, 7909560, y, 1005840, 320040, tk, size=12.5, bold=True, color=INK)
    d.text(s, 8915400, y, 2331720, 320040, lab, size=11, color=MUTED)

# ── 12b · MONTE CARLO ───────────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "FEATURE · CERTAINTY", "One line implies a certainty that doesn't exist",
         eyecolor=PURPLE,
         sub="A retirement projection is a distribution, not a forecast. So show the "
             "distribution.")
d.bullets(s, [
    "10,000 simulated paths of household value — median plus 10th / 25th / 75th / 90th "
    "percentile bands, on a KPI tile that opens a fan chart",
    "User controls for horizon, annual contribution and return assumptions",
    "Seeded and reproducible — a projection that changes every time you open it is a slot "
    "machine, not a plan",
], y=2103120, w=6126480, h=2011680, size=15)
d.rect(s, M, 4297680, 6126480, 1005840, fill=SOFT, line=LINE, radius=0.05)
d.text(s, M + 228600, 4434840, 5669280, 731520, [
    [("No new data source.  ", {"bold": True, "color": PURPLE}),
     ("Volatility is estimated from the 52-week high/low already sitting in every row of the "
      "table. Look hard at what your data can already tell you before you go shopping for a "
      "new feed.", {})],
], size=12, color=INK)
# mock fan
d.rect(s, 7406640, 2103120, 4114800, 3108960, fill=CARD, line=LINE, radius=0.05)
d.text(s, 7680960, 2331720, 3566160, 320040, "Median @ year 10", size=13, bold=True, color=INK)
_N = 14
for i in range(_N):
    f = i / float(_N - 1)
    cx = 7680960 + i * 241578
    cy = 4300000 - int(800000 * f)                 # rising median
    ho = int(40000 + 600000 * f)                   # 10th-90th half-height
    hi = int(ho * 0.45)                            # 25th-75th half-height
    d.rect(s, cx, cy - ho, 200000, ho * 2, fill=LILAC)
    d.rect(s, cx, cy - hi, 200000, hi * 2, fill=PURPLE)
    d.rect(s, cx, cy - 12000, 200000, 24000, fill=INK)
d.text(s, 7680960, 4700000, 3566160, 240000,
       "median path · bands = 25th\u201375th and 10th\u201390th percentile", size=9.5, color=MUTED)

# ── 12c · DCF VALUATION ─────────────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "FEATURE · VALUATION", "Is it actually worth what it costs?", eyecolor=GREEN,
         sub="The question the whole dashboard implies — and never answered until now.")
d.bullets(s, [
    "A three-stage DCF per held stock — analyst consensus, a 10-year fade, then a rebuilt "
    "steady state — shown as a fair-value strip beside the market price",
    "Every judgement input is a user control: capex intensity, fade, terminal growth, beta, "
    "risk-free rate, equity risk premium, consensus haircut",
    "A reverse DCF reports the discount rate the market is implicitly applying to consensus",
    "A Trefis-style segment sum-of-the-parts tab values each business line separately",
], y=2011680, w=6126480, h=2286000, size=15)
d.rect(s, M, 4389120, 6126480, 1005840, fill=CARD, line=LINE, radius=0.05)
d.text(s, M + 228600, 4526280, 5669280, 731520, [
    [("Benchmarked externally, or it doesn't ship.  ", {"bold": True, "color": GREEN}),
     ("This model was built three times before it was right — each time \u2018verified\u2019 against "
      "the vendor's own model, which only ever proved a flawed model had been faithfully "
      "reproduced. Internal consistency is not correctness.", {})],
], size=12, color=INK)
# mock valuation strip
d.rect(s, 7406640, 2011680, 4114800, 3200400, fill=CARD, line=LINE, radius=0.05)
d.text(s, 7680960, 2240280, 3566160, 320040, "Fair value", size=13, bold=True, color=INK)
_rows = [("MSFT", "$583", "vs $600 independent", GREEN),
         ("NVDA", "$267", "vs $280 independent", GREEN),
         ("AMD", "segment view refused", "segments overstate revenue 12%", AMBER)]
for i, (tk, val, note, col) in enumerate(_rows):
    y = 2743200 + i * 731520
    d.rect(s, 7680960, y, 3383280, 594360, fill=SOFT, line=LINE, radius=0.05)
    d.text(s, 7817160, y + 68580, 1097280, 274320, tk, size=12, bold=True, color=INK)
    d.text(s, 8778240, y + 68580, 2194560, 274320, val, size=12, bold=True, color=col)
    d.text(s, 7817160, y + 320040, 3109280, 240000, note, size=9.5, color=MUTED)
d.text(s, 7680960, 4983480, 3566160, 240000,
       "refusing to render, with a reason, is a feature", size=9.5, color=MUTED)

# ── 12d · PORTFOLIO LOOK-THROUGH ────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "FEATURE · X-RAY", "What do I actually own?", eyecolor=RED,
         sub="Single-name risk arrives through funds, not just through shares — and the "
             "account tabs cannot see it.")
d.bullets(s, [
    "Every held ETF and mutual fund decomposed into its underlying companies, then combined "
    "with the directly held shares",
    "A duplicate-tracker check — three separate S&P 500 trackers held side by side is one bet, "
    "not three",
    "Explicit buckets for each fund's long tail, cash inside funds, and any fund not yet "
    "mapped — reconciled to the household total, to the dollar",
], y=2103120, w=6126480, h=2011680, size=15)
d.rect(s, M, 4297680, 6126480, 1005840, fill=SOFT, line=LINE, radius=0.05)
d.text(s, M + 228600, 4434840, 5669280, 731520, [
    [("A look-through that hides its gaps  ", {"bold": True, "color": RED}),
     ("understates the exact concentration it exists to reveal. So while any fund is unmapped, "
      "every figure is presented as a floor \u2014 \u2018at least\u2019, never a total.", {})],
], size=12, color=INK)
# mock exposure bars
d.rect(s, 7406640, 2103120, 4114800, 3108960, fill=CARD, line=LINE, radius=0.05)
d.text(s, 7680960, 2331720, 3566160, 320040, "True single-name exposure", size=13, bold=True,
       color=INK)
_ex = [("NVDA", 0.00, 1.00, "6.2%"), ("AAPL", 0.55, 0.45, "5.1%"),
       ("MSFT", 0.40, 0.60, "4.8%"), ("AMZN", 0.00, 0.80, "3.4%"),
       ("AVGO", 0.30, 0.35, "2.2%")]
for i, (tk, dpart, fpart, pct) in enumerate(_ex):
    y = 2834640 + i * 457200
    d.text(s, 7680960, y, 914400, 274320, tk, size=11, bold=True, color=INK)
    bx, bwid = 8686800, 1828800
    if dpart > 0:
        d.rect(s, bx, y + 45720, int(bwid * dpart * 0.5), 137160, fill=GREEN)
    d.rect(s, bx + int(bwid * dpart * 0.5), y + 45720, int(bwid * fpart * 0.5), 137160,
           fill=PURPLE)
    d.text(s, 10744200, y, 640080, 274320, pct, size=10.5, bold=True, color=INK)
d.text(s, 7680960, 4937760, 3566160, 240000,
       "green = shares you hold · purple = reaching you through funds", size=9.5, color=MUTED)

# ── 13 · FUNDAMENTALS POPUP ─────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "FEATURE · FUNDAMENTALS POPUP", "Click any ticker — get the story", eyecolor=PURPLE)
d.bullets(s, [
    "Factsheet — a plain-English summary + a vitals grid (sector, HQ, founded, employees, "
    "market cap / issuer, expense, AUM, holdings)",
    "A data-derived “Did you know?” fun fact on every name",
    "An 8-factor risk summary (incl. supply-chain) and an advisor scorecard",
    "Cached locally, refreshed daily — fast on repeat opens",
], y=1828800, w=6126480, h=2926080, size=15)
d.rect(s, 7406640, 1828800, 4114800, 3200400, fill=CARD, line=LINE, radius=0.05)
d.text(s, 7680960, 2057400, 3566160, 365760, "Apple Inc.", size=16, bold=True, color=INK)
d.text(s, 7680960, 2395728, 3566160, 274320,
       "Technology · Consumer Electronics · Cupertino", size=10.5, color=MUTED)
d.text(s, 7680960, 2697480, 3566160, 822960,
       "Designs iPhone, Mac & iPad plus a fast-growing Services arm; among the world's most "
       "valuable companies.", size=11.5, color=INK)
d.rect(s, 7680960, 3611880, 3566160, 868680, fill=LILAC)
d.text(s, 7863840, 3703320, 3246120, 731520, [
    [("Did you know?  ", {"bold": True, "color": PURPLE}),
     ("Founded in a garage in 1976 — now one of a handful of companies ever worth over "
      "$3 trillion.", {})],
], size=11, color=INK)

# ── 14 · RULES-BASED RANKING ────────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "FEATURE · RULES-BASED RANKING", "Rankings you can customize — and trust",
         eyecolor=BLUE)
d.text(s, M, 1371600, CONTENT_W, 640080, [
    [("•  Sector, fund/category, and growth-core rankings on a declarative rules engine", {})],
    [("•  Tune weights, thresholds, and IF/THEN conditions live", {})],
    [("•  Every score traces to its rules — a per-row breakdown popover explains it", {})],
], size=14, color=INK, space=6)
d.image(s, "08_ranking_engine.png", M, 2560320, CONTENT_W, 3566160)

# ── 15 · DAILY AUTOMATION ───────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "FEATURE · DAILY AUTOMATION", "Zero-touch, every morning")
d.text(s, M, 1371600, CONTENT_W, 457200,
       "Fetch → build → guard → import → reprice → recompute → alert → project income → "
       "regenerate → publish — with a safe fallback at every rung.",
       size=14, color=MUTED)
d.image(s, "04_scheduled_workflow.png", M, 1965960, CONTENT_W, 4160520)

# ── 16 · RELIABILITY ────────────────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "FEATURE · RELIABILITY & DATA INTEGRITY", "Built to be trusted", eyecolor=GREEN)
d.tiles(s, [
    ("image-13-1.png", GREEN, "Integrity guard",
     "Every row checked against its 52-week band; corrupt / #N/A values are suppressed and "
     "flagged with a ⚠ badge — and suppressed rows can never raise a false alert."),
    ("image-13-2.png", GREEN, "Fallback ladder",
     "FMP live → workbook tail → carry-forward. A total outage degrades safely to the last "
     "good data."),
    ("image-13-3.png", GREEN, "Never blanks",
     "Holdings, rankings, charts, history, alerts and the dividend calendar are all preserved "
     "on any failure — the dashboard never breaks."),
])

# ── 17 · OBSERVABILITY ──────────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "OBSERVABILITY", "Logging & a debug view", eyecolor=PURPLE)
d.tiles(s, [
    ("image-14-1.png", PURPLE, "Structured run logs",
     "Every daily run writes a per-run event stream, a human-readable report, and an "
     "append-only history — auto-pruned. A failed run is diagnosable after the fact."),
    ("image-14-2.png", PURPLE, "Opt-in debug panel",
     "Enable with ?debug=1 or Shift+D (off by default). See connector calls, data provenance, "
     "cache state, and errors — with a one-click copy-report for support."),
], top=2103120, h=2926080)

# ── 18 · UNDER THE HOOD ─────────────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "UNDER THE HOOD", "System architecture")
d.image(s, "01_architecture.png", M, 1554480, CONTENT_W, 4663440)

# ── 19 · WHY IT HOLDS UP ────────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "WHY IT HOLDS UP", "Engineering discipline the AI applied")
d.tiles(s, [
    ("image-16-1.png", BLUE, "Modular build",
     "A single-file app assembled from 23 concern-based modules — with a byte-for-byte "
     "golden-master guard."),
    ("image-16-2.png", BLUE, "Verified changes",
     "Edits validated in a stubbed-DOM harness, not just a syntax check — regressions caught "
     "before publish."),
    ("image-16-3.png", BLUE, "Provider adapter",
     "All connector I/O behind one boundary — a provider or shape change is a one-place edit."),
    ("image-16-4.png", BLUE, "Integrity gates",
     "Guards and safe fallbacks at every automated step — speed without fragility."),
])

# ── 20 · THE METHOD ─────────────────────────────────────────────────────────────
s = d.slide(NAVY)
d.text(s, M, 457200, 8229600, 274320, "THE AI-ADOPTION METHOD", size=12, bold=True, color=GOLD)
d.text(s, M, 777240, CONTENT_W, 731520, "How it was built",
       size=32, bold=True, color=CARD, font=DISPLAY)
steps = [("1", "Clarify", GOLD, "image-17-1.png",
          "A couple of pointed questions before each step — build the right thing, not the wrong thing."),
         ("2", "Plan + mock", GREEN, "image-17-2.png",
          "Show a plan and a visual mockup, agree, then build — no surprises."),
         ("3", "Small steps", BLUE, "image-17-3.png",
          "One capability per prompt, validated, then the next — reversible and traceable."),
         ("4", "Harden + doc", RED, "image-17-4.png",
          "Guards, tests, and living docs so it's production-grade, not a demo.")]
for i, (num, head, col, img, body) in enumerate(steps):
    x = M + i * 2788920
    d.rect(s, x, 1920240, 2560320, 3200400, fill=DARKC, line=DARKL, radius=0.05)
    d.icon(s, x + 960120, 2240280, 640080, col, img)
    d.text(s, x + 182880, 3063240, 2194560, 365760, num, size=14, bold=True,
           color=STEEL, font=DISPLAY)
    d.text(s, x + 182880, 3383280, 2194560, 365760, head, size=16, bold=True, color=CARD)
    d.text(s, x + 228600, 3794760, 2103120, 1188720, body, size=11.5, color=LIGHT)
d.footer(s, dark=True)

# ── 21 · DELIVERY TIMELINE ──────────────────────────────────────────────────────
s = d.slide(CARD)
d.chrome(s, "DELIVERY TIMELINE", f"Baseline v{VERSION} is the starting point")
d.rect(s, 1737360, 3108960, 8714232, 12700, fill=LINE, line=LINE)
miles = [("v1.0", "Baseline", BLUE,
          "9-tab dashboard, live portfolios, retirement, ranking, valuation, look-through, daily automation, docs"),
         ("Next", "Backlog", MUTED,
          "Tax-lot, PDF/email, corporate actions, saved scenarios, security/CI, Morningstar FVE")]
step = 8714232 // (len(miles) - 1)
for i, (v, lab, col, body) in enumerate(miles):
    cx = 1737360 + i * step
    d.rect(s, cx - 146304, 2962656, 292608, 292608, fill=col, radius=0.5)
    d.text(s, cx - 822960, 2148840, 1645920, 365760, v, size=20, bold=True,
           color=col, font=DISPLAY, align=PP_ALIGN.CENTER)
    d.text(s, cx - 822960, 2542032, 1645920, 320040, lab, size=13, bold=True,
           color=INK, align=PP_ALIGN.CENTER)
    _bw = step - 91440
    d.text(s, cx - _bw // 2, 3429000, _bw, 1463040, body, size=10, color=MUTED,
           align=PP_ALIGN.CENTER)

# ── 22 · BUSINESS VALUE ─────────────────────────────────────────────────────────
s = d.slide(NAVY)
d.text(s, M, 457200, 8229600, 274320, "BUSINESS VALUE", size=12, bold=True, color=GOLD)
d.text(s, M, 777240, CONTENT_W, 731520, "What it's worth",
       size=32, bold=True, color=CARD, font=DISPLAY)
vals = [("Hours → 0", GOLD, "weekly upkeep replaced by a zero-touch daily refresh"),
        ("100% local", GREEN, "holdings never leave the machine — private by design"),
        ("1 file", BLUE, "shareable, installs nothing, works offline"),
        ("Adaptable", RED, "rules, universe, and features change by conversation")]
gap = 274320
w = (CONTENT_W - gap * 3) // 4
for i, (big, col, cap) in enumerate(vals):
    x = M + i * (w + gap)
    d.rect(s, x, 2103120, w, 2560320, fill=DARKC, line=DARKL, radius=0.06)
    d.text(s, x + 228600, 2469936, w - 457200, 731520, big, size=24, bold=True,
           color=col, font=DISPLAY)
    d.text(s, x + 228600, 3291840, w - 457200, 1097280, cap, size=12, color=LIGHT)
d.footer(s, dark=True)

# ── 23 · FOR CLIENT ENGAGEMENTS ─────────────────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "FOR YOUR CLIENT ENGAGEMENTS", "The repeatable pattern")
d.tiles(s, [
    ("image-20-1.png", GOLD, "Find the workflow",
     "A spreadsheet-and-manual-lookup process that's stale, fragmented, and high-effort."),
    ("image-20-2.png", GREEN, "Build conversationally",
     "Plan, mock, and build feature by feature — in the client's own language."),
    ("image-20-3.png", BLUE, "Harden",
     "Guards, safe fallbacks, and tests so it runs unattended and never breaks."),
    ("image-20-4.png", PURPLE, "Document & automate",
     "Living docs, a scheduled refresh, and observability — hand off a system, not a script."),
])

# ── 24 · BOB 2.0 — AI AGENT INFRASTRUCTURE ─────────────────────────────────────
s = d.slide(NAVY)
d.text(s, M, 457200, 8229600, 274320, "AI AGENT INFRASTRUCTURE · BOB 2.0",
       size=12, bold=True, color=GOLD)
d.text(s, M, 777240, CONTENT_W, 731520,
       "MCP, Markdown & the agent configuration layer",
       size=32, bold=True, color=CARD, font=DISPLAY)

# Three dark tiles: MCP transport, markdown files, env / secrets
_bob_tiles = [
    (GOLD,  "Remote HTTP MCP",
     "Both data connectors are fully-managed remote HTTP servers — no npm install, no local "
     "process to babysit. Bigdata connects automatically. FMP binds :8080 on startup.",
     [("Bigdata", "https://mcp.bigdata.com/", GOLD),
      ("FMP",     "http://localhost:8080/mcp", AMBER)]),
    (TEAL,  "Markdown instruction files",
     "Agent behaviour lives in plain Markdown: AGENTS.md (invariants, pipeline rules, "
     "provider constraints), daily_refresh.md (9-step playbook), and the .bob/ rules-ask layer.",
     [("AGENTS.md", "invariants · pipeline · data-provider rules", TEAL),
      ("daily_refresh.md", "10-step daily playbook", TEAL)]),
    (PURPLE, "Env vars & secrets",
     "API keys never live in code. FMP_ACCESS_TOKEN and BIGDATA_API_KEY are read from .env "
     "(gitignored). mcp.json carries only the transport URL and auto-approved tool list.",
     [("FMP_ACCESS_TOKEN", ".env (gitignored)", PURPLE),
      ("BIGDATA_API_KEY", ".bob/mcp.json → x-api-key header", PURPLE)]),
]
tile_gap = 274320
tile_w = (CONTENT_W - tile_gap * 2) // 3
tile_top = 1828800
tile_h = 4023360

for ti, (col, head, body, rows) in enumerate(_bob_tiles):
    tx = M + ti * (tile_w + tile_gap)
    d.rect(s, tx, tile_top, tile_w, tile_h, fill=DARKC, line=DARKL, radius=0.05)
    # colour bar top
    d.rect(s, tx, tile_top, tile_w, 45720, fill=col)
    # icon dot
    d.rect(s, tx + 228600, tile_top + 182880, 137160, 137160, fill=col, radius=0.5)
    d.text(s, tx + 228600, tile_top + 365760, tile_w - 457200, 411480,
           head, size=15, bold=True, color=CARD)
    d.text(s, tx + 228600, tile_top + 822960, tile_w - 457200, 1188720,
           body, size=11, color=LIGHT)
    # mini spec rows
    for ri, (lab, val, rcol) in enumerate(rows):
        ry = tile_top + 2194560 + ri * 731520
        d.rect(s, tx + 182880, ry, tile_w - 365760, 594360, fill="0A1D35", line=DARKL, radius=0.04)
        d.text(s, tx + 365760, ry + 91440, tile_w - 640080, 228600,
               lab, size=10.5, bold=True, color=col)
        d.text(s, tx + 365760, ry + 320040, tile_w - 640080, 228600,
               val, size=9.5, color=SLATE)

d.footer(s, dark=True)

# ── 25 · CODE ANALYSIS — INIT & MIGRATION ───────────────────────────────────────
s = d.slide(SOFT)
d.chrome(s, "CODE ANALYSIS · INIT & MIGRATION",
         "Understanding and evolving an existing codebase",
         eyecolor=BLUE,
         sub="Bob 2.0 can read an entire repo on first open, reason about its structure, "
             "and guide migrations — not just write net-new code.")

# Left column: bullets
d.bullets(s, [
    "GetSymbolsOverview and FindSymbol map every class, function, and module across "
    "the project in seconds — no manual spelunking",
    "FindReferencingSymbols traces all callers of any symbol, so a refactor's blast "
    "radius is known before the first edit",
    "Init flow: Bob reads AGENTS.md, inspects src/dash/, scripts/, tests/ and builds "
    "a mental model of the codebase before answering any question",
    "Migration: replace a provider, rename a constant, restructure a module — Bob "
    "proposes the diff, verifies it compiles in the stubbed-DOM harness, then ships both "
    "HTML copies in a single agent turn",
], y=2194560, w=6126480, h=3474720, size=14)

# Right column: two stacked cards
card_x = 7040880
card_w = 4557384

# Card 1 — Init
d.rect(s, card_x, 1828800, card_w, 2011680, fill=CARD, line=LINE, radius=0.05)
d.rect(s, card_x, 1828800, card_w, 45720, fill=BLUE)
d.text(s, card_x + 228600, 1920240, card_w - 457200, 320040,
       "Init  ·  first open", size=13, bold=True, color=BLUE)
_init = [("Read", "AGENTS.md + TASKS.md"),
         ("Map", "src/dash/ · scripts/ · tests/ · data/"),
         ("Verify", "golden-master + test suite green"),
         ("Ready", "can answer any question about the codebase")]
for ri, (step, detail) in enumerate(_init):
    ry = 2286000 + ri * 365760
    d.text(s, card_x + 228600, ry, 685800, 274320, step,
           size=10.5, bold=True, color=BLUE)
    d.text(s, card_x + 1005840, ry, card_w - 1234440, 274320, detail,
           size=10.5, color=INK)

# Card 2 — Migration
d.rect(s, card_x, 4023360, card_w, 2011680, fill=CARD, line=LINE, radius=0.05)
d.rect(s, card_x, 4023360, card_w, 45720, fill=PURPLE)
d.text(s, card_x + 228600, 4114800, card_w - 457200, 320040,
       "Migration  ·  provider swap example", size=13, bold=True, color=PURPLE)
_mig = [("Identify", "FindReferencingSymbols on the old provider"),
        ("Diff", "propose targeted apply_diff per file"),
        ("Test", "node --test + pytest in the same turn"),
        ("Deploy", "assemble + inject + write both HTML copies")]
for ri, (step, detail) in enumerate(_mig):
    ry = 4480560 + ri * 365760
    d.text(s, card_x + 228600, ry, 685800, 274320, step,
           size=10.5, bold=True, color=PURPLE)
    d.text(s, card_x + 1005840, ry, card_w - 1234440, 274320, detail,
           size=10.5, color=INK)

# ── 26 · CLOSING ────────────────────────────────────────────────────────────────
s = d.slide(NAVY)
for i, col in enumerate((GREEN, AMBER, RED)):
    d.rect(s, M + i * 384048, 2331720, 310896, 310896, fill=col, radius=0.2)
d.text(s, M, 2834640, 10789920, 822960, "From spreadsheet to self-updating system —",
       size=34, bold=True, color=CARD, font=DISPLAY)
d.text(s, M, 3611880, 10789920, 822960, "conversationally.",
       size=34, bold=True, color=GOLD, font=DISPLAY)
d.text(s, M, 4754880, 10789920, 365760,
       f"Tech Drawdown · Version {VERSION} · Rules-based analytics, not financial advice.",
       size=13, color=SLATE)

d.save()
